"""
SPS MBB Dipole — 200 GeV vs 26 GeV Comparison — PowerPoint Presentation
========================================================================
Uses the CERN Corporate 16:9 template (2024).
Generates all analysis plots and assembles them into a PPTX presentation.
Covers: B1/b2/b3 comparison, transfer function, eddy-current settling,
        SFTPRO history dependence, per-supercycle analysis, outlier removal.
"""

import re
import sys
from pathlib import Path

import numpy as np
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# ── Project imports ──────────────────────────────────────────────────
sys.path.insert(0, str(Path(__file__).resolve().parent))
from rotating_coil_analyzer.analysis.kn_pipeline import load_segment_kn_txt
from rotating_coil_analyzer.analysis.utility_functions import (
    compute_block_averaged_range,
    detect_plateau_turns,
    classify_current,
    find_contiguous_groups,
    process_kn_pipeline,
    build_harmonic_rows,
    mad_sigma_clip,
    compute_level_stats,
    diff_sigma,
)
from rotating_coil_analyzer.ingest.channel_detect import robust_range
from rotating_coil_analyzer.presentation.pptx_helpers import (
    LY_TITLE_LOGO, LY_LAST,
    CERN_BLUE, DARK_GREY, WHITE, LIGHT_BG,
    set_footer, add_picture_to_slide,
    slide_title_only, slide_chapter, slide_bullets,
    slide_text_and_pic, slide_table, slide_two_images,
    savefig,
)

# ── Global plot style ────────────────────────────────────────────────
plt.rcParams.update({
    "figure.dpi": 180,
    "axes.grid": True,
    "grid.alpha": 0.3,
    "font.size": 11,
    "axes.titlesize": 13,
    "axes.labelsize": 11,
})

# =====================================================================
# CONFIGURATION
# =====================================================================
BASE_DIR = Path(r"C:\Users\albellel\python-projects\rotating-coil-analyzer\measurements\2026_02_06")

SESSION_200 = BASE_DIR / "01_200_extended" / "20260206_144537_SPS_MBB"
SUBDIR_200 = "20260206_144559_MBB"

SESSION_26 = BASE_DIR / "03_26_extended" / "20260206_151808_SPS_MBB"
SUBDIR_26 = "20260206_151827_MBB"

SEGMENT = "NCS"
KN_PATH = Path(r"C:\Users\albellel\python-projects\rotating-coil-analyzer\measurements\20251212_171026_SPS_MBA\CRMMMMH_AV-00000001\Kn_values_Seg_Main_A_AC.txt")

TEMPLATE_PATH = Path(r"C:\Users\albellel\python-projects\rotating-coil-analyzer\presentation_template\CERN Corporate_16x9 PPT_v2024_converted.pptx")

MAGNET_ORDER = 1
R_REF = 0.02
SAMPLES_PER_TURN = 1024
Ns = SAMPLES_PER_TURN
OPTIONS = ("dri", "rot", "cel", "fed")
DRIFT_MODE = "legacy"
MIN_B1_T = 1e-4
MAX_ZR = 0.01
N_BLOCKS = 10
PLATEAU_I_RANGE_MAX = 2.5

N_LAST_TURNS_INJ = 18
N_LAST_TURNS_HIGH = None
N_SIGMA_CLIP = 5
FLIP_FIELD_SIGN = False

OUT_DIR = Path(r"C:\Users\albellel\python-projects\rotating-coil-analyzer\output")
OUT_DIR.mkdir(exist_ok=True)
IMG_DIR = OUT_DIR / "pptx_images_sps_mbb"
IMG_DIR.mkdir(exist_ok=True)

FILE_PAT = re.compile(
    r"Run_(\d+)_I_([\d.]+)A_(N?CS)_raw_measurement_data\.txt$"
)

SPS_CURRENT_THRESHOLDS = {
    "zero": 50, "pre-ramp": 200, "injection": 500,
    "flat-low": 2000, "flat-mid": 4000,
}
ANALYSIS_LABELS = {"injection", "flat-mid", "flat-high"}

FOOTER_TEXT = "Alberto Bellelli | SPS MBB Dipole 200 vs 26 GeV | MMM Section"
DATE_TEXT = "12 February 2026"
_F = dict(footer_date=DATE_TEXT, footer_text=FOOTER_TEXT)


# =====================================================================
# DATA PROCESSING
# =====================================================================

def load_and_process(session_dir, meas_subdir, kn, dataset_label=""):
    run_dir = session_dir / meas_subdir
    ncs_files = []
    for f in sorted(run_dir.iterdir()):
        match = FILE_PAT.search(f.name)
        if match and match.group(3) == SEGMENT:
            ncs_files.append(f)
    assert ncs_files, f"No {SEGMENT} raw files found in {run_dir}"
    raw_file = ncs_files[0]
    print(f"\n  Dataset: {dataset_label}")
    print(f"  Raw file: {raw_file.name}")

    raw = np.loadtxt(raw_file)
    n_turns = raw.shape[0] // Ns
    n_keep = n_turns * Ns
    print(f"  {n_turns} turns, {raw[-1,0] - raw[0,0]:.1f} s")

    t_all = raw[:n_keep, 0].reshape(n_turns, Ns)
    flux_abs_all = raw[:n_keep, 1].reshape(n_turns, Ns)
    flux_cmp_all = raw[:n_keep, 2].reshape(n_turns, Ns)
    I_all = raw[:n_keep, 3].reshape(n_turns, Ns)

    I_mean_quick = I_all.mean(axis=1)
    best_turn = np.argmax(np.abs(I_mean_quick))
    r1 = robust_range(raw[best_turn*Ns:(best_turn+1)*Ns, 1])
    r2 = robust_range(raw[best_turn*Ns:(best_turn+1)*Ns, 2])
    if r2 > r1:
        flux_abs_all = raw[:n_keep, 2].reshape(n_turns, Ns)
        flux_cmp_all = raw[:n_keep, 1].reshape(n_turns, Ns)
        print("  (flux columns swapped)")

    I_mean = I_all.mean(axis=1)
    t_mean = t_all.mean(axis=1)
    I_range, I_blocks = compute_block_averaged_range(I_all, Ns, N_BLOCKS)

    plateau_info = detect_plateau_turns(I_blocks, I_mean, I_range, PLATEAU_I_RANGE_MAX)
    is_plateau = plateau_info["is_plateau"]

    turn_label = np.array(["ramp"] * n_turns, dtype=object)
    for i in range(n_turns):
        if is_plateau[i]:
            turn_label[i] = classify_current(I_mean[i])

    is_analysis = np.array([l in ANALYSIS_LABELS for l in turn_label])
    plateau_indices = np.where(is_analysis)[0]
    print(f"  {len(plateau_indices)} plateau turns")

    if len(plateau_indices) == 0:
        empty = pd.DataFrame()
        return empty, empty, [], [], t_mean, I_mean, turn_label, is_plateau

    idx = plateau_indices
    result, C_merged, C_units, ok_main = process_kn_pipeline(
        flux_abs_turns=flux_abs_all[idx],
        flux_cmp_turns=flux_cmp_all[idx],
        t_turns=t_all[idx],
        I_turns=I_all[idx],
        kn=kn, r_ref=R_REF, magnet_order=MAGNET_ORDER,
        options=OPTIONS, drift_mode=DRIFT_MODE,
        min_b1_T=MIN_B1_T, max_zr=MAX_ZR,
    )

    extra = [
        {"global_turn": int(idx[t]), "label": str(turn_label[idx[t]]),
         "I_range_A": float(I_range[idx[t]])}
        for t in range(len(idx))
    ]
    rows = build_harmonic_rows(result, C_merged, C_units, ok_main, MAGNET_ORDER, extra)
    df = pd.DataFrame(rows)

    if FLIP_FIELD_SIGN:
        t_cols = [c for c in df.columns if c.endswith("_T")]
        df[t_cols] *= -1

    # Group injection by supercycle
    inj_mask_global = (turn_label == "injection")
    sc_groups_inj = find_contiguous_groups(inj_mask_global, min_length=2)
    df["sc_idx"] = -1
    settled_idx = []

    for gi, (gs, ge) in enumerate(sc_groups_inj):
        group_globals = set(range(gs, ge + 1))
        gmask = df["global_turn"].isin(group_globals) & (df["label"] == "injection")
        df.loc[gmask, "sc_idx"] = gi
        group_rows = df.index[gmask]
        if N_LAST_TURNS_INJ is not None and len(group_rows) > N_LAST_TURNS_INJ:
            settled_idx.extend(group_rows[-N_LAST_TURNS_INJ:])
        else:
            settled_idx.extend(group_rows)

    # Group flat-high by supercycle
    fh_mask_global = (turn_label == "flat-high")
    sc_groups_fh = find_contiguous_groups(fh_mask_global, min_length=2)
    for gi, (gs, ge) in enumerate(sc_groups_fh):
        group_globals = set(range(gs, ge + 1))
        gmask = df["global_turn"].isin(group_globals) & (df["label"] == "flat-high")
        df.loc[gmask, "sc_idx"] = gi
        group_rows = df.index[gmask]
        if N_LAST_TURNS_HIGH is not None and len(group_rows) > N_LAST_TURNS_HIGH:
            settled_idx.extend(group_rows[-N_LAST_TURNS_HIGH:])
        else:
            settled_idx.extend(group_rows)

    df_settled = df.loc[sorted(settled_idx)].copy()

    # MAD sigma clipping
    n_before = len(df_settled)
    df_settled, clip_removed = mad_sigma_clip(df_settled, "B1_T", N_SIGMA_CLIP)
    n_clipped = n_before - len(df_settled)
    if n_clipped > 0:
        print(f"  Sigma clip: removed {n_clipped} turns ({clip_removed})")

    n_inj = len(df_settled[df_settled["label"] == "injection"])
    n_fh = len(df_settled[df_settled["label"] == "flat-high"])
    print(f"  Final settled: {n_inj} injection, {n_fh} flat-high")

    return df, df_settled, sc_groups_inj, sc_groups_fh, t_mean, I_mean, turn_label, is_plateau


# =====================================================================
# MAIN
# =====================================================================
def main():
    print("=" * 70)
    print("  SPS MBB 200 GeV vs 26 GeV -- PowerPoint Presentation Generator")
    print("=" * 70)

    # ── Load Kn ──
    print("\n[1] Loading Kn calibration...")
    kn = load_segment_kn_txt(str(KN_PATH))
    print(f"  {len(kn.orders)} harmonics from {KN_PATH.name}")

    # ── Process both datasets ──
    print("\n[2] Processing 200 GeV Extended...")
    df_200, dfs_200, sc_inj_200, sc_fh_200, t_200, I_200, lbl_200, plat_200 = \
        load_and_process(SESSION_200, SUBDIR_200, kn, "200 GeV Extended")

    print("\n[3] Processing 26 GeV Extended...")
    df_26, dfs_26, sc_inj_26, sc_fh_26, t_26, I_26, lbl_26, plat_26 = \
        load_and_process(SESSION_26, SUBDIR_26, kn, "26 GeV Extended")

    # ── Color scheme ──
    COL_200 = "tab:blue"
    COL_26 = "tab:orange"

    # =================================================================
    # GENERATE ALL PLOTS
    # =================================================================
    print("\n[4] Generating plots...")

    # ── P01: Current profile side by side ──
    print("  P01: Current profiles")
    label_colors = {"injection": "tab:green", "flat-mid": "tab:purple",
                    "flat-high": "tab:blue"}
    fig, axes = plt.subplots(1, 2, figsize=(14, 5), sharey=True)
    for ax, t_m, I_m, lbl, title in [
        (axes[0], t_200, I_200, lbl_200, "200 GeV Extended"),
        (axes[1], t_26, I_26, lbl_26, "26 GeV Extended"),
    ]:
        ax.plot(t_m, I_m, ".-", markersize=1, linewidth=0.3,
                color="lightgrey", zorder=0)
        for lab, col in label_colors.items():
            mask = lbl == lab
            idx = np.where(mask)[0]
            if len(idx) > 0:
                ax.scatter(t_m[idx], I_m[idx], s=6, color=col, zorder=2, label=lab)
        ax.set_xlabel("Time (s)")
        ax.set_title(title)
        ax.legend(fontsize=8, loc="upper right")
    axes[0].set_ylabel("I (A)")
    fig.suptitle("Current Profile -- 200 GeV vs 26 GeV Extended", fontsize=13, y=1.02)
    plt.tight_layout()
    img_current = savefig(fig, IMG_DIR,"P01_current_profile")

    # ── P02: Eddy-current settling at injection ──
    print("  P02: Eddy-current settling")
    fig, axes = plt.subplots(1, 2, figsize=(14, 6))
    for ax, df_all, sc_grps, ds_name in [
        (axes[0], df_200, sc_inj_200, "200 GeV"),
        (axes[1], df_26, sc_inj_26, "26 GeV"),
    ]:
        inj = df_all[(df_all["label"] == "injection") & df_all["ok_main"]].copy()
        if len(inj) == 0:
            continue
        n_sc = inj["sc_idx"].max() + 1
        for sc in range(n_sc):
            sc_data = inj[inj["sc_idx"] == sc].sort_values("global_turn")
            if len(sc_data) < 2:
                continue
            turn_in_sc = np.arange(len(sc_data))
            if N_LAST_TURNS_INJ is not None and len(sc_data) > N_LAST_TURNS_INJ:
                ref = sc_data["B1_T"].values[-N_LAST_TURNS_INJ:].mean()
            else:
                ref = sc_data["B1_T"].mean()
            delta = (sc_data["B1_T"].values - ref) * 1e6
            ax.plot(turn_in_sc, delta, ".-", markersize=2, alpha=0.35,
                    color=plt.cm.viridis(sc / max(n_sc - 1, 1)), linewidth=0.8)
        if N_LAST_TURNS_INJ is not None:
            typical_len = int(inj.groupby("sc_idx").size().median())
            cutoff = typical_len - N_LAST_TURNS_INJ
            if cutoff > 0:
                ax.axvline(cutoff - 0.5, color="red", linestyle="--",
                           linewidth=1.5, label=f"Cutoff (skip first {cutoff})")
                ax.axvspan(-0.5, cutoff - 0.5, color="red", alpha=0.06)
        ax.axhline(0, color="grey", linewidth=0.5)
        ax.set_xlabel("Turn within supercycle")
        ax.set_ylabel("$\\Delta$B1 from settled mean ($\\mu$T)")
        ax.set_title(f"{ds_name} -- {n_sc} supercycles")
        ax.legend(fontsize=8)
    fig.suptitle(f"Eddy-Current Settling at Injection (last {N_LAST_TURNS_INJ} turns selected)",
                 fontsize=12, y=1.04)
    plt.tight_layout()
    img_settling = savefig(fig, IMG_DIR,"P02_settling")

    # ── P03: TF at injection per supercycle ──
    print("  P03: Transfer function at injection")
    fig, axes = plt.subplots(1, 2, figsize=(14, 5))
    ax = axes[0]
    for ds_name, dfs, col in [("200 GeV", dfs_200, COL_200),
                               ("26 GeV", dfs_26, COL_26)]:
        inj = dfs[(dfs["label"] == "injection") & dfs["ok_main"]].copy()
        if len(inj) == 0:
            continue
        inj["TF"] = inj["B1_T"] / (inj["I_mean_A"] / 1000.0)
        sc_tf = inj.groupby("sc_idx")["TF"].agg(["mean", "std"]).reset_index()
        ax.errorbar(sc_tf["sc_idx"], sc_tf["mean"], yerr=sc_tf["std"],
                    fmt="o-", markersize=4, capsize=2, color=col, alpha=0.8,
                    label=ds_name)
    ax.set_xlabel("Supercycle index")
    ax.set_ylabel("TF = B1/I (T/kA)")
    ax.set_title("TF vs Supercycle (injection, settled)")
    ax.legend()

    ax = axes[1]
    box_data, box_labels, box_colors = [], [], []
    for ds_name, dfs, col in [("200 GeV", dfs_200, COL_200),
                               ("26 GeV", dfs_26, COL_26)]:
        inj = dfs[(dfs["label"] == "injection") & dfs["ok_main"]].copy()
        if len(inj) == 0:
            continue
        tf = (inj["B1_T"] / (inj["I_mean_A"] / 1000.0)).values
        box_data.append(tf)
        box_labels.append(f"{ds_name}\n(N={len(tf)})")
        box_colors.append(col)
    if box_data:
        bp = ax.boxplot(box_data, tick_labels=box_labels, patch_artist=True)
        for patch, c in zip(bp["boxes"], box_colors):
            patch.set_facecolor(c)
            patch.set_alpha(0.5)
    ax.set_ylabel("TF = B1/I (T/kA)")
    ax.set_title("TF Distribution (injection, settled)")
    fig.suptitle("Transfer Function at Injection (~301 A)", fontsize=13, y=1.02)
    plt.tight_layout()
    img_tf_inj = savefig(fig, IMG_DIR,"P03_tf_injection")

    # ── P04: Per-supercycle injection harmonics ──
    print("  P04: Per-supercycle injection harmonics")
    fig, axes = plt.subplots(1, 3, figsize=(16, 5))
    harmonics = [("B1_T", "B1 (T)"), ("b2_units", "b2 (units)"), ("b3_units", "b3 (units)")]
    for ax, (col_name, ylabel) in zip(axes, harmonics):
        for ds_name, dfs, col in [("200 GeV", dfs_200, COL_200),
                                   ("26 GeV", dfs_26, COL_26)]:
            inj = dfs[(dfs["label"] == "injection") & dfs["ok_main"]]
            if len(inj) == 0:
                continue
            sc_avg = inj.groupby("sc_idx")[col_name].agg(["mean", "std"]).reset_index()
            ax.errorbar(sc_avg["sc_idx"], sc_avg["mean"], yerr=sc_avg["std"],
                        fmt="o-", markersize=4, capsize=2, color=col, alpha=0.8,
                        label=ds_name)
        ax.set_xlabel("Supercycle index")
        ax.set_ylabel(ylabel)
        ax.legend(fontsize=9)
    axes[0].set_title("B1 per supercycle")
    axes[1].set_title("b2 per supercycle")
    axes[2].set_title("b3 per supercycle")
    fig.suptitle(f"Per-Supercycle Injection Harmonics (last {N_LAST_TURNS_INJ} settled turns)",
                 fontsize=13, y=1.02)
    plt.tight_layout()
    img_sc_inj = savefig(fig, IMG_DIR,"P04_sc_injection")

    # ── P05: SFTPRO per-supercycle (B1, b2, b3, TF) ──
    print("  P05: SFTPRO per-supercycle")
    fig, axes = plt.subplots(2, 2, figsize=(14, 10))
    harmonics_sftpro = [
        ("B1_T", "B1 (T)"), ("b2_units", "b2 (units)"),
        ("b3_units", "b3 (units)"), ("TF_TperkA", "TF (T/kA)"),
    ]
    for ax, (col_name, ylabel) in zip(axes.ravel(), harmonics_sftpro):
        for ds_name, dfs, col in [("After 200 GeV MD1", dfs_200, COL_200),
                                   ("After 26 GeV MD1", dfs_26, COL_26)]:
            fh = dfs[(dfs["label"] == "flat-high") & dfs["ok_main"]].copy()
            if len(fh) == 0:
                continue
            if col_name == "TF_TperkA":
                fh["TF_TperkA"] = fh["B1_T"] / (fh["I_mean_A"] / 1000.0)
            sc_avg = fh.groupby("sc_idx")[col_name].agg(["mean", "std"]).reset_index()
            ax.errorbar(sc_avg["sc_idx"], sc_avg["mean"], yerr=sc_avg["std"],
                        fmt="o-", markersize=4, capsize=2, color=col, alpha=0.8,
                        label=ds_name)
        ax.set_xlabel("Supercycle index")
        ax.set_ylabel(ylabel)
        ax.legend(fontsize=9)
    axes[0, 0].set_title("B1 per supercycle (SFTPRO)")
    axes[0, 1].set_title("b2 per supercycle (SFTPRO)")
    axes[1, 0].set_title("b3 per supercycle (SFTPRO)")
    axes[1, 1].set_title("TF per supercycle (SFTPRO)")
    fig.suptitle("SFTPRO Flat-Top (~4815 A): After 200 GeV MD1 vs After 26 GeV MD1",
                 fontsize=13, y=1.02)
    plt.tight_layout()
    img_sc_sftpro = savefig(fig, IMG_DIR,"P05_sc_sftpro")

    # ── P06: B1 comparison (injection + SFTPRO) ──
    print("  P06: B1 comparison")
    fig, axes = plt.subplots(1, 2, figsize=(14, 5))
    for ax, lab, title in [
        (axes[0], "injection", f"Injection (~301 A, last {N_LAST_TURNS_INJ}/SC)"),
        (axes[1], "flat-high", "Top of SFTPRO (~4815 A)"),
    ]:
        for ds_name, dfs, col in [("200 GeV", dfs_200, COL_200),
                                   ("26 GeV", dfs_26, COL_26)]:
            sub = dfs[(dfs["label"] == lab) & dfs["ok_main"]]
            if len(sub) == 0:
                continue
            ax.plot(sub["time_s"].values, sub["B1_T"].values,
                    ".", markersize=3, alpha=0.6, color=col, label=ds_name)
        ax.set_xlabel("Time (s)")
        ax.set_ylabel("B1 (T)")
        ax.set_title(title)
        ax.legend(fontsize=9)
    fig.suptitle("B1 Main Field Comparison (settled turns)", fontsize=13, y=1.02)
    plt.tight_layout()
    img_b1 = savefig(fig, IMG_DIR,"P06_b1_comparison")

    # ── P07: b2 comparison ──
    print("  P07: b2 comparison")
    fig, axes = plt.subplots(1, 2, figsize=(14, 5))
    for ax, lab, title in [
        (axes[0], "injection", f"Injection (~301 A, last {N_LAST_TURNS_INJ}/SC)"),
        (axes[1], "flat-high", "Top of SFTPRO (~4815 A)"),
    ]:
        for ds_name, dfs, col in [("200 GeV", dfs_200, COL_200),
                                   ("26 GeV", dfs_26, COL_26)]:
            sub = dfs[(dfs["label"] == lab) & dfs["ok_main"]]
            if len(sub) == 0:
                continue
            ax.plot(sub["time_s"].values, sub["b2_units"].values,
                    ".", markersize=3, alpha=0.6, color=col, label=ds_name)
        ax.axhline(0, color="grey", linewidth=0.5)
        ax.set_xlabel("Time (s)")
        ax.set_ylabel("b2 (units)")
        ax.set_title(title)
        ax.legend(fontsize=9)
    fig.suptitle("b2 Normal Quadrupole Comparison (settled turns)", fontsize=13, y=1.02)
    plt.tight_layout()
    img_b2 = savefig(fig, IMG_DIR,"P07_b2_comparison")

    # ── P08: b3 comparison ──
    print("  P08: b3 comparison")
    fig, axes = plt.subplots(1, 2, figsize=(14, 5))
    for ax, lab, title in [
        (axes[0], "injection", f"Injection (~301 A, last {N_LAST_TURNS_INJ}/SC)"),
        (axes[1], "flat-high", "Top of SFTPRO (~4815 A)"),
    ]:
        for ds_name, dfs, col in [("200 GeV", dfs_200, COL_200),
                                   ("26 GeV", dfs_26, COL_26)]:
            sub = dfs[(dfs["label"] == lab) & dfs["ok_main"]]
            if len(sub) == 0:
                continue
            ax.plot(sub["time_s"].values, sub["b3_units"].values,
                    ".", markersize=3, alpha=0.6, color=col, label=ds_name)
        ax.axhline(0, color="grey", linewidth=0.5)
        ax.set_xlabel("Time (s)")
        ax.set_ylabel("b3 (units)")
        ax.set_title(title)
        ax.legend(fontsize=9)
    fig.suptitle("b3 Normal Sextupole Comparison (settled turns)", fontsize=13, y=1.02)
    plt.tight_layout()
    img_b3 = savefig(fig, IMG_DIR,"P08_b3_comparison")

    # ── P09: Box plots (2x2) ──
    print("  P09: Box plots")
    fig, axes = plt.subplots(2, 2, figsize=(14, 10))
    axes_flat = axes.ravel()
    plot_specs = [
        ("B1_T", "B1 (T)", "B1 Main Field"),
        ("b2_units", "b2 (units)", "b2 Normal Quadrupole"),
        ("b3_units", "b3 (units)", "b3 Normal Sextupole"),
        ("TF_TperkA", "TF (T/kA)", "Transfer Function"),
    ]
    for ax, (col_name, ylabel, title) in zip(axes_flat, plot_specs):
        box_data, box_labels, box_colors = [], [], []
        for ds_name, dfs, base_col in [("200 GeV", dfs_200, COL_200),
                                        ("26 GeV", dfs_26, COL_26)]:
            for lab, short in [("injection", "Inj"), ("flat-high", "SFTPRO")]:
                sub = dfs[(dfs["label"] == lab) & dfs["ok_main"]].copy()
                if len(sub) == 0:
                    continue
                if col_name == "TF_TperkA":
                    vals = (sub["B1_T"] / (sub["I_mean_A"] / 1000.0)).values
                else:
                    vals = sub[col_name].values
                box_data.append(vals)
                box_labels.append(f"{ds_name}\n{short}\n(N={len(sub)})")
                box_colors.append(base_col)
        if box_data:
            bp = ax.boxplot(box_data, tick_labels=box_labels, patch_artist=True)
            for patch, c in zip(bp["boxes"], box_colors):
                patch.set_facecolor(c)
                patch.set_alpha(0.5)
        if col_name not in ("B1_T", "TF_TperkA"):
            ax.axhline(0, color="grey", linewidth=0.5)
        ax.set_ylabel(ylabel)
        ax.set_title(title)
        ax.tick_params(axis="x", labelsize=8)
    fig.suptitle("Distribution Comparison: 200 GeV vs 26 GeV (settled turns)",
                 fontsize=13, y=1.02)
    plt.tight_layout()
    img_boxplots = savefig(fig, IMG_DIR,"P09_boxplots")

    # ── P10: Bar chart ──
    print("  P10: Bar chart")
    fig, axes = plt.subplots(2, 2, figsize=(14, 10))
    axes_flat = axes.ravel()
    datasets = [("200 GeV", dfs_200), ("26 GeV", dfs_26)]
    op_points = [("injection", "Injection (MD1)"), ("flat-high", "Top SFTPRO")]
    colors_ds = {"200 GeV": COL_200, "26 GeV": COL_26}
    bar_harmonics = [
        ("B1_T", "B1 (T)"), ("b2_units", "b2 (units)"),
        ("b3_units", "b3 (units)"), ("TF_TperkA", "TF (T/kA)"),
    ]
    x = np.arange(len(op_points))
    width = 0.35
    for ax, (col_name, ylabel) in zip(axes_flat, bar_harmonics):
        for i, (ds_name, dfs) in enumerate(datasets):
            means, stds = [], []
            for lab, _ in op_points:
                sub = dfs[(dfs["label"] == lab) & dfs["ok_main"]].copy()
                if len(sub) > 0:
                    if col_name == "TF_TperkA":
                        vals = sub["B1_T"] / (sub["I_mean_A"] / 1000.0)
                    else:
                        vals = sub[col_name]
                    means.append(vals.mean())
                    stds.append(vals.std())
                else:
                    means.append(0)
                    stds.append(0)
            offset = (i - 0.5) * width
            ax.bar(x + offset, means, width, yerr=stds, capsize=4,
                   color=colors_ds[ds_name], alpha=0.7, label=ds_name)
        ax.set_xticks(x)
        ax.set_xticklabels([desc for _, desc in op_points], fontsize=9)
        ax.set_ylabel(ylabel)
        ax.set_title(col_name.replace("_", " "))
        ax.legend(fontsize=9)
    fig.suptitle(f"Mean +/- Std: 200 GeV vs 26 GeV (settled, last {N_LAST_TURNS_INJ}/SC at inj.)",
                 fontsize=13, y=1.02)
    plt.tight_layout()
    img_barchart = savefig(fig, IMG_DIR,"P10_barchart")

    # =================================================================
    # COMPUTE SUMMARY STATISTICS FOR TABLES
    # =================================================================
    print("\n[5] Computing summary statistics...")

    s200_inj = compute_level_stats(dfs_200, "injection")
    s200_fh = compute_level_stats(dfs_200, "flat-high")
    s26_inj = compute_level_stats(dfs_26, "injection")
    s26_fh = compute_level_stats(dfs_26, "flat-high")

    # =================================================================
    # BUILD POWERPOINT
    # =================================================================
    print("\n[6] Building PowerPoint with CERN template...")
    prs = Presentation(str(TEMPLATE_PATH))

    # ── TITLE SLIDE ──
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_LOGO])
    slide.placeholders[0].text = "SPS MBB Dipole\n200 GeV vs 26 GeV Comparison"
    if 1 in {ph.placeholder_format.idx for ph in slide.placeholders}:
        slide.placeholders[1].text = (
            "Field Quality | Transfer Function | History Dependence | Eddy Currents\n"
            "Alberto Bellelli -- CERN MMM Section -- February 2026")

    # ── SECTION 1: MEASUREMENT SETUP ──
    slide_chapter(prs, "Measurement Setup",
                  "SPS MBB dipole -- two cycle configurations on the same magnet", **_F)

    slide_bullets(prs, "Measurement Overview", [
        "Magnet: SPS MBB (Main Bending dipole), warm iron-dominated",
        "Date: 2026-02-06, two sessions ~30 min apart",
        "Segment: NCS (non-connection side)",
        "Kn: cross-session calibration from MBA (15 harmonics)",
        "Supercycle: LHC_pilot -> MD1 -> SFTPRO, repeated 20 times",
        "Two datasets: 200 GeV extended (MD1 at higher I) vs 26 GeV extended (MD1 at lower I)",
        "Total: ~1061-1064 turns per dataset (~17.7 min each)",
    ], **_F)

    slide_table(prs, "Dataset Summary",
        ["Parameter", "200 GeV Extended", "26 GeV Extended"],
        [["Session", "01_200_extended", "03_26_extended"],
         ["Timestamp", "2026-02-06 14:45", "2026-02-06 15:18"],
         ["Total turns", "1061", "1064"],
         ["Injection turns", f"{s200_inj['N']}", f"{s26_inj['N']}"],
         ["SFTPRO turns", f"{s200_fh['N']}", f"{s26_fh['N']}"],
         ["Injection I", f"{s200_inj['I_mean']:.1f} A", f"{s26_inj['I_mean']:.1f} A"],
         ["SFTPRO I", f"{s200_fh['I_mean']:.1f} A", f"{s26_fh['I_mean']:.1f} A"],
         ["Supercycles", "20", "20"]],
        col_widths=[3.0, 4.0, 4.0], **_F)

    slide_bullets(prs, "Analysis Pipeline", [
        "Pipeline: Bottura standard (MTA-IN-97-007)",
        "Corrections: drift + rotation + centre loc. + feed-down",
        "Merge: abs_upto_m_cmp_above (B1 from absolute, n>1 from compensated)",
        "Reference radius: R_ref = 20 mm",
        f"Plateau detection: block-averaged I range < {PLATEAU_I_RANGE_MAX} A",
        f"Settling: last {N_LAST_TURNS_INJ} of ~24 injection turns/SC (skip first ~6)",
        f"Outlier removal: MAD sigma clipping at {N_SIGMA_CLIP} sigma on B1",
        "FLIP_FIELD_SIGN = False (raw sign convention)",
    ], **_F)

    # ── SECTION 2: CURRENT PROFILE & SETTLING ──
    slide_chapter(prs, "Current Profile & Eddy-Current Settling",
                  "Supercycle structure and settling correction", **_F)

    slide_title_only(prs, "Current Profile -- 200 GeV vs 26 GeV", img_current, **_F)

    slide_title_only(prs, "Eddy-Current Settling at Injection", img_settling,
        notes=f"Each injection plateau ~24 turns. First ~6 contaminated by eddy currents. "
              f"Keep last {N_LAST_TURNS_INJ} turns per supercycle.", **_F)

    slide_bullets(prs, "Settling Correction & Outlier Removal", [
        f"Injection: ~24 turns per supercycle, first ~6 show eddy-current settling",
        f"Keep last {N_LAST_TURNS_INJ} turns per supercycle (skip settling transient)",
        f"SFTPRO: all turns kept (no settling concern at cycle top)",
        f"Outlier removal: MAD sigma clipping (N = {N_SIGMA_CLIP} sigma) on B1",
        "Catches ramp-boundary turns (last turn of group where ramp has started)",
        f"200 GeV: {s200_fh['N']} SFTPRO turns after cleaning",
        f"26 GeV: {s26_fh['N']} SFTPRO turns after cleaning",
        "Note: 26 GeV injection has 14x larger B1 scatter (71 vs 5 uT std)",
    ], **_F)

    # ── SECTION 3: INJECTION COMPARISON ──
    slide_chapter(prs, "Injection Comparison (~301 A)",
                  "Same within-supercycle ramp pattern, different full history", **_F)

    slide_title_only(prs, "Transfer Function at Injection (~301 A)", img_tf_inj, **_F)

    slide_title_only(prs, "Per-Supercycle Injection: B1, b2, b3", img_sc_inj, **_F)

    slide_title_only(prs, "B1 Main Field -- Injection vs SFTPRO", img_b1, **_F)

    # ── SECTION 4: SFTPRO HISTORY DEPENDENCE ──
    slide_chapter(prs, "SFTPRO Flat-Top: History Dependence",
                  "Does the preceding MD1 cycle affect the field at ~4815 A?", **_F)

    slide_bullets(prs, "SFTPRO History Dependence -- Context", [
        "SFTPRO flat-top (~4815 A) preceded by different MD1 cycles:",
        "   200 GeV dataset: SFTPRO follows a 200 GeV MD1 (higher flat-top I)",
        "   26 GeV dataset: SFTPRO follows a 26 GeV MD1 (lower flat-top I)",
        "If iron magnetisation history affects the field at SFTPRO,",
        "   B1 and harmonics should differ between datasets",
        "Caveat: the full history differs too (separate sessions, ~30 min apart)",
        "   Cannot isolate single-cycle from multi-cycle memory",
        "Iron hysteresis has memory extending many cycles back",
    ], **_F)

    slide_title_only(prs, "SFTPRO Per-Supercycle: B1, b2, b3, TF", img_sc_sftpro, **_F)

    # SFTPRO difference table
    dB1, dB1_err, dB1_sig = diff_sigma(s200_fh, s26_fh, "B1")
    db2, db2_err, db2_sig = diff_sigma(s200_fh, s26_fh, "b2")
    db3, db3_err, db3_sig = diff_sigma(s200_fh, s26_fh, "b3")
    dTF, dTF_err, dTF_sig = diff_sigma(s200_fh, s26_fh, "TF")

    slide_table(prs, "SFTPRO Flat-Top: 200 GeV vs 26 GeV Difference",
        ["Quantity", "200 GeV", "26 GeV", "Delta", "Significance"],
        [["B1 (T)", f"{s200_fh['B1_mean']:.6f}", f"{s26_fh['B1_mean']:.6f}",
          f"{dB1*1e6:+.0f} uT", f"{dB1_sig:.1f} sigma"],
         ["b2 (units)", f"{s200_fh['b2_mean']:+.4f}", f"{s26_fh['b2_mean']:+.4f}",
          f"{db2:+.4f}", f"{db2_sig:.1f} sigma"],
         ["b3 (units)", f"{s200_fh['b3_mean']:+.4f}", f"{s26_fh['b3_mean']:+.4f}",
          f"{db3:+.4f}", f"{db3_sig:.1f} sigma"],
         ["TF (T/kA)", f"{s200_fh['TF_mean']:.5f}", f"{s26_fh['TF_mean']:.5f}",
          f"{dTF*1e3:+.2f} mT/kA", f"{dTF_sig:.1f} sigma"]],
        col_widths=[2.0, 2.5, 2.5, 2.5, 2.0], **_F)

    slide_bullets(prs, "SFTPRO -- Transfer Function: 2.8 Sigma", [
        "TF = B1/I is the most sensitive indicator (2.8 sigma)",
        "Dividing by per-turn current removes current-correlated B1 scatter",
        "Raw B1 std ~0.47 mT includes current spread effect",
        "TF std ~0.10 mT/kA -- much tighter distribution",
        "The ~0.03 mT/kA TF offset approaches the 3-sigma threshold",
        "Direct signature of history-dependent iron permeability at 4815 A",
        "Consistent across all 20 supercycles (systematic offset, not drift)",
        "But: separate sessions -- cannot exclude session-to-session systematics",
    ], **_F)

    # ── SECTION 5: HARMONIC COMPARISON ──
    slide_chapter(prs, "Harmonic Comparison: b2 & b3",
                  "Quadrupole and sextupole field errors", **_F)

    slide_two_images(prs, "b2 and b3 Comparison (Injection + SFTPRO)",
                     img_b2, img_b3, "b2 (quadrupole)", "b3 (sextupole)", **_F)

    slide_bullets(prs, 'What Does "X Sigma" Mean?', [
        "sigma = |difference| / uncertainty on the difference",
        "   uncertainty = sqrt(std\u00b2/N dataset1 + std\u00b2/N dataset2)",
        "< 2 sigma: no evidence of a real difference (could be noise)",
        "2\u20133 sigma: suggestive but not conclusive",
        "> 3 sigma: strong evidence the difference is real (< 0.3% chance of noise)",
        "IMPORTANT: high sigma \u2260 large difference!",
        "   B1 at injection: ~22 sigma but only 7 \u00b5T (60 ppm) -- real but tiny",
        "   TF at SFTPRO: 2.8 sigma -- borderline, probably real but not conclusive",
    ], **_F)

    # Injection difference table
    dB1_i, dB1_i_err, dB1_i_sig = diff_sigma(s200_inj, s26_inj, "B1")
    db2_i, db2_i_err, db2_i_sig = diff_sigma(s200_inj, s26_inj, "b2")
    db3_i, db3_i_err, db3_i_sig = diff_sigma(s200_inj, s26_inj, "b3")
    dTF_i, dTF_i_err, dTF_i_sig = diff_sigma(s200_inj, s26_inj, "TF")

    slide_table(prs, "Full Difference Table (200 GeV - 26 GeV)",
        ["Quantity", "Injection delta", "Inj. sigma", "SFTPRO delta", "SFTPRO sigma"],
        [["B1", f"{dB1_i*1e6:+.0f} uT", f"{dB1_i_sig:.1f}",
          f"{dB1*1e6:+.0f} uT", f"{dB1_sig:.1f}"],
         ["b2 (units)", f"{db2_i:+.4f}", f"{db2_i_sig:.1f}",
          f"{db2:+.4f}", f"{db2_sig:.1f}"],
         ["b3 (units)", f"{db3_i:+.4f}", f"{db3_i_sig:.1f}",
          f"{db3:+.4f}", f"{db3_sig:.1f}"],
         ["TF (T/kA)", f"{dTF_i*1e3:+.2f} mT/kA", f"{dTF_i_sig:.1f}",
          f"{dTF*1e3:+.2f} mT/kA", f"{dTF_sig:.1f}"]],
        col_widths=[2.0, 2.5, 2.0, 2.5, 2.0], **_F)

    # ── SECTION 6: BOX PLOTS & BAR CHART ──
    slide_chapter(prs, "Statistical Distributions",
                  "Box plots and bar charts for all quantities", **_F)

    slide_title_only(prs, "Box Plots: B1, b2, b3, TF", img_boxplots, **_F)

    slide_title_only(prs, "Bar Chart: Mean +/- Std", img_barchart, **_F)

    # ── SECTION 7: SUMMARY TABLE ──
    slide_chapter(prs, "Summary & Conclusions",
                  "Key results and measurement limitations", **_F)

    slide_table(prs, "Summary Statistics (Settled + Cleaned Turns)",
        ["Dataset", "Op. Point", "N", "B1 (T)", "b2 (units)", "b3 (units)", "TF (T/kA)"],
        [["200 GeV", "Injection", f"{s200_inj['N']}",
          f"{s200_inj['B1_mean']:.6f}", f"{s200_inj['b2_mean']:+.3f}",
          f"{s200_inj['b3_mean']:+.3f}", f"{s200_inj['TF_mean']:.4f}"],
         ["200 GeV", "SFTPRO", f"{s200_fh['N']}",
          f"{s200_fh['B1_mean']:.6f}", f"{s200_fh['b2_mean']:+.3f}",
          f"{s200_fh['b3_mean']:+.3f}", f"{s200_fh['TF_mean']:.4f}"],
         ["26 GeV", "Injection", f"{s26_inj['N']}",
          f"{s26_inj['B1_mean']:.6f}", f"{s26_inj['b2_mean']:+.3f}",
          f"{s26_inj['b3_mean']:+.3f}", f"{s26_inj['TF_mean']:.4f}"],
         ["26 GeV", "SFTPRO", f"{s26_fh['N']}",
          f"{s26_fh['B1_mean']:.6f}", f"{s26_fh['b2_mean']:+.3f}",
          f"{s26_fh['b3_mean']:+.3f}", f"{s26_fh['TF_mean']:.4f}"]],
        col_widths=[1.5, 1.5, 1.0, 2.0, 1.8, 1.8, 1.8], **_F)

    slide_bullets(prs, "Conclusions: Can We Distinguish History Effects?", [
        "Hysteresis (DC): marginal -- TF at SFTPRO 2.8 sigma, B1 1.7 sigma",
        "   Suggestive but below 3 sigma; cannot exclude session systematics",
        "Eddy currents: visible within each dataset (settling in first ~6 turns)",
        "   26 GeV has 14x larger injection scatter -- different settling dynamics",
        "b3 sextupole: ~4.7 sigma difference at injection (only >3 sigma finding)",
        "   Cannot attribute to specific mechanism (history vs systematic)",
        "TF saturation: 3.1% drop from 301 A to 4815 A -- identical in both datasets",
        "Limitation: separate sessions, cannot isolate single-cycle from multi-cycle memory",
    ], **_F)

    slide_bullets(prs, "Recommendations for Future Measurements", [
        "Interleave cycles: alternate 200 GeV and 26 GeV MD1 within same session",
        "   Eliminates session-to-session offsets for direct comparison",
        "Pre-cycle: standardize initial magnetisation state",
        "Monitor temperature: rule out thermal contributions to b3 difference",
        "Longer flat-top: more SFTPRO turns per SC would improve statistics",
        "   Currently 3-4 turns/SC at SFTPRO -- limited by cycle structure",
        "Cross-session Kn: continue using MBA calibration (in-session Kn are zeros)",
    ], **_F)

    # ── LAST SLIDE ──
    slide = prs.slides.add_slide(prs.slide_layouts[LY_LAST])
    txBox = slide.shapes.add_textbox(Inches(3), Inches(3), Inches(7), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank you"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = CERN_BLUE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "alberto.bellelli@cern.ch"
    p2.font.size = Pt(16)
    p2.font.color.rgb = DARK_GREY
    p2.alignment = PP_ALIGN.CENTER

    # ── Save ──
    pptx_path = OUT_DIR / "SPS_MBB_200GeV_vs_26GeV_Comparison.pptx"
    prs.save(str(pptx_path))
    n_plots = len(list(IMG_DIR.glob("*.png")))
    print(f"\n{'=' * 70}")
    print(f"  Presentation saved: {pptx_path}")
    print(f"  {len(prs.slides)} slides, {n_plots} plots generated")
    print(f"{'=' * 70}")


if __name__ == "__main__":
    main()
