#!/usr/bin/env python
"""Build a focused MC62 Test 03 presentation from executed notebooks.

Extracts PNG images from the already-executed Test 03 analysis and
eddy-current notebooks, then assembles a ~25-slide PPTX using the
CERN template.

Usage
-----
    python build_test03_presentation.py           # build PPTX
    python build_test03_presentation.py --count   # just count images

Prerequisites
-------------
- python-pptx, Pillow, lxml
- Notebooks executed in-place (via nbconvert --execute --inplace)
"""
from __future__ import annotations

import sys
from pathlib import Path

# -- Reuse helpers from the campaign presentation builder ------------------
from build_presentation import (
    extract_images,
    add_title_slide,
    add_chapter_slide,
    add_image_slide,
    add_text_slide,
    add_last_slide,
    _safe,
)

# -- slide_table from the shared pptx_helpers module -----------------------
sys.path.insert(0, str(Path(__file__).resolve().parents[2] / "presentation"))
from pptx_helpers import slide_table, slide_bullets

from pptx import Presentation
from lxml import etree

# -- Paths -----------------------------------------------------------------
REPO = Path(r"C:\Users\albellel\python-projects\rotating-coil-analyzer")
NB_DIR = REPO / "rotating_coil_analyzer" / "notebooks" / "LEAR_MC62"
TEMPLATE_PPTX = NB_DIR / "MC62_2Hz_staircase_presentation.pptx"
OUTPUT_PPTX = NB_DIR / "MC62_test03_presentation.pptx"

ANALYSIS_NB = NB_DIR / "analysis" / "2026-02-16_03_staircase_2Hz.ipynb"
EDDY_NB = NB_DIR / "eddy_current" / "2026-02-16_03_staircase_2Hz.ipynb"


# -- Count mode ------------------------------------------------------------

def count_images():
    """Print image counts for both Test 03 notebooks."""
    for label, path in [("analysis_03", ANALYSIS_NB), ("eddy_03", EDDY_NB)]:
        if path.exists():
            imgs = extract_images(path)
            print(f"  {label}: {len(imgs)} images")
        else:
            print(f"  {label}: NOT FOUND at {path}")


# -- Build -----------------------------------------------------------------

def build():
    """Build the focused Test 03 presentation."""
    # Load template and clear existing slides
    prs = Presentation(str(TEMPLATE_PPTX))
    while len(prs.slides._sldIdLst) > 0:
        sldId = prs.slides._sldIdLst[0]
        rId = sldId.get(etree.QName(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
            "id"))
        if rId:
            prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(sldId)

    # Extract images
    print("Extracting images from executed notebooks...")
    a03 = extract_images(ANALYSIS_NB) if ANALYSIS_NB.exists() else []
    e03 = extract_images(EDDY_NB) if EDDY_NB.exists() else []
    print(f"  analysis_03: {len(a03)} images")
    print(f"  eddy_03:     {len(e03)} images")

    # ==================================================================
    # Slide 1: Title
    # ==================================================================
    add_title_slide(
        prs,
        "MC62 Test 03 \u2014 2 Hz Staircase",
        "LEAR C-shaped Dipole\n"
        "Feb 16, 2026 \u2014 Afternoon Measurement\n"
        "A. Bellelli \u2014 CERN",
    )

    # ==================================================================
    # Slide 2: Setup
    # ==================================================================
    add_text_slide(prs, "Measurement Setup & Analysis Pipeline",
        "2 Hz staircase (120 rpm), 512 samples/turn, ~740 turns/plateau\n"
        "Current cycle: 0\u2192+200\u21920\u2192\u2212200\u21920 A, "
        "20 A steps, 1 A/s ramp\n"
        "No shims, streaming binary format\n\n"
        "Pipeline: dri + rot (cel/fed auto-disabled \u2014 UNSAFE for dipole)\n"
        "R_ref = 33.0 mm, external Kn calibration\n"
        "Averaging: N_LAST = 340 turns (after eddy-current settling)\n\n"
        "PCBs: Integral (R45, 30 harmonics) + Central (DQ, 15 harmonics)\n"
        "Eddy-current fit: single-exp B1(t) = B_inf + A exp(\u2212t/\u03c4)")

    # ==================================================================
    # Chapter 1: Measurement Quality
    # ==================================================================
    add_chapter_slide(prs, 1,
        "Measurement Quality",
        "Plateau detection, turn classification, FDI diagnostics")

    # Slide 4: Current profile
    if _safe(a03, 1):
        add_image_slide(prs, "Current Profile & Plateau Detection",
                        a03[1],
                        "41 plateaus detected, 20 A steps, "
                        "block-averaged range threshold")

    # Slide 5: Turn classification map
    if _safe(a03, 2):
        add_image_slide(prs, "Turn Classification Map",
                        a03[2],
                        "Green = plateau, Orange = ramp, Grey = precycle")

    # Slide 6: FDI diagnostic summary
    add_text_slide(prs, "FDI Diagnostic Summary",
        "Stuck-channel detection: 40/40 transitions OK\n\n"
        "All FDI channels show proper transitions between current levels.\n"
        "No stuck ADC channels detected in any plateau.\n\n"
        "Comparison with Test 04:\n"
        "  \u2022 Test 03: 0/40 stuck transitions (clean)\n"
        "  \u2022 Test 04: 5/40 stuck transitions\n\n"
        "\u2192 Test 03 is the reference 2 Hz dataset.")

    # ==================================================================
    # Chapter 2: Harmonic Results
    # ==================================================================
    add_chapter_slide(prs, 2,
        "Harmonic Results",
        "Hysteresis curves, multipole spectrum, FFMM validation")

    # Slide 8: Zoomed 2x2 hysteresis (img 11)
    if _safe(a03, 11):
        add_image_slide(prs, "Harmonic Hysteresis (I \u2260 0)",
                        a03[11],
                        "B1, b2, b3, TF vs current \u2014 "
                        "b2 \u2248 \u221216, b3 \u2248 \u221212 units")

    # Slide 9: B1 full range (img 6)
    if _safe(a03, 6):
        add_image_slide(prs, "B1 vs Current (Full Range)",
                        a03[6])

    # Slide 10: b2 hysteresis (img 7)
    if _safe(a03, 7):
        add_image_slide(prs, "b2 (Quadrupole) Hysteresis",
                        a03[7],
                        "b2 \u2248 \u221216 units (C-shape asymmetry, no shims)")

    # Slide 11: b3 hysteresis (img 8)
    if _safe(a03, 8):
        add_image_slide(prs, "b3 (Sextupole) Hysteresis",
                        a03[8])

    # Slide 12: Multipole spectrum (img 12)
    if _safe(a03, 12):
        add_image_slide(prs, "Multipole Spectrum at Peak Current",
                        a03[12])

    # Slide 13: FFMM parity (img 13)
    if _safe(a03, 13):
        add_image_slide(prs, "FFMM C++ Parity Check",
                        a03[13],
                        "B_main max |diff| = 1.81e\u221213 T (machine precision)")

    # ==================================================================
    # Chapter 3: Eddy Current Analysis
    # ==================================================================
    add_chapter_slide(prs, 3,
        "Eddy Current Analysis",
        "Single-exp model: B1(t) = B_inf + A exp(\u2212t/\u03c4)\n"
        "2-pass fit with 5\u03c3 MAD outlier clipping")

    # Slide 15: Settling curves
    if _safe(e03, 0):
        add_image_slide(prs, "Settling Curves (B1 vs Turn)",
                        e03[0])

    # Slide 16: Fits grid
    if _safe(e03, 1):
        add_image_slide(prs, "Exponential Fits \u2014 All Plateaus",
                        e03[1])

    # Slide 17: Tau vs current
    if _safe(e03, 2):
        add_image_slide(prs, "\u03c4 vs Current",
                        e03[2],
                        "\u03c4 = 24.2 \u00b1 7.3 s (GOOD fits); "
                        "\u03bc_r dependence visible")

    # Slide 18: Settling bias
    if _safe(e03, 3):
        add_image_slide(prs, "Settling Bias (B1, b2, b3)",
                        e03[3])

    # Slide 19: N_LAST sensitivity
    if _safe(e03, 4):
        add_image_slide(prs, "N_LAST_TURNS Sensitivity Study",
                        e03[4],
                        "b3 bias < 0.02 units for all N_LAST values")

    # Slide 20: Double-exp comparison
    if _safe(e03, 6):
        add_image_slide(prs, "Single vs Double Exponential Comparison",
                        e03[6],
                        "Marginal R\u00b2 improvement \u2014 "
                        "single-exp adequate for this magnet")

    # ==================================================================
    # Chapter 4: Data Quality & Conclusions
    # ==================================================================
    add_chapter_slide(prs, 4,
        "Data Quality\n& Conclusions",
        "Fit quality breakdown, physical explanation, key findings")

    # Slide 22: Data quality table
    slide_table(prs, "Eddy Current Fit Quality Summary",
        headers=["Category", "Count", "Details"],
        rows=[
            ["Total fitted",  "38",              "Runs with |I| \u2265 10 A"],
            ["GOOD",          "18",              "R\u00b2 \u2265 0.9"],
            ["MARGINAL",      "18",              "R\u00b2 < 0.9 (physical)"],
            ["WEAK_SIGNAL",   "2",               "Runs 16, 18"],
            ["FIT_FAILED",    "0",               "\u2014"],
            ["FDI stuck",     "0 / 40",          "All transitions OK"],
            ["\u03c4 (GOOD)", "24.2 \u00b1 7.3 s", "Range 8.1\u201340.0 s"],
            ["Outliers/run",  "0\u201321",       "MAD 5\u03c3 clip"],
        ],
        col_widths=[3.2, 2.5, 5.8],
    )

    # Slide 23: Physical explanation
    slide_bullets(prs, "Why R\u00b2 < 0.9 Is Not a Data Quality Problem",
        [
            "At high current (|I| \u2265 180 A): iron saturates \u2192 "
            "\u03bc_r drops",
            "Lower \u03bc_r: smaller eddy-current amplitude A",
            "Smaller A relative to noise: lower R\u00b2 "
            "(model still physically correct)",
            "Descending branch: permeability hysteresis further "
            "reduces transient amplitude",
            "WEAK_SIGNAL (runs 16, 18): |A| < 3\u00d7 noise at "
            "+80 A, +40 A descending",
            "The 18 GOOD fits provide reliable \u03c4 values "
            "(\u03c4 = 24.2 \u00b1 7.3 s)",
        ],
    )

    # Slide 24: Key findings
    slide_bullets(prs, "Key Findings \u2014 Test 03",
        [
            "B1 at 200 A: 0.2185 T (integral)",
            "TF at 200 A: 1.093 T/kA",
            "b2: \u221216 units (C-shape asymmetry, no shims)",
            "b3: \u221212 units (stable)",
            "Eddy current \u03c4: 8\u201340 s (GOOD fits), "
            "\u03bc_r-dependent",
            "N_LAST = 340 turns safely excludes 5\u00d7\u03c4_max",
            "FFMM parity: machine-precision "
            "(< 2\u00d710\u207b\u00b9\u00b3 T)",
            "FDI: all 40 transitions clean \u2192 "
            "reference 2 Hz dataset",
        ],
    )

    # -- Last slide --
    add_last_slide(prs)

    # Save
    prs.save(str(OUTPUT_PPTX))
    print(f"\nPresentation saved to: {OUTPUT_PPTX}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    if "--count" in sys.argv:
        print("Image counts for Test 03 notebooks:")
        count_images()
    else:
        build()
