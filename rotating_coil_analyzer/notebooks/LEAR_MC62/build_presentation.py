#!/usr/bin/env python
"""Build the MC62 measurement campaign presentation from executed notebooks.

Extracts PNG images from Jupyter notebooks that have already been
executed in-place (with outputs embedded), then assembles a PPTX
using the CERN template.

Usage
-----
    python build_presentation.py           # build PPTX
    python build_presentation.py --count   # just count images per notebook

Prerequisites
-------------
- python-pptx, Pillow
- Notebooks executed in-place (via nbconvert --execute --inplace)
"""
from __future__ import annotations

import base64
import io
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN

# -- Paths -----------------------------------------------------------------
REPO = Path(r"C:\Users\albellel\python-projects\rotating-coil-analyzer")
NB_DIR = REPO / "rotating_coil_analyzer" / "notebooks" / "LEAR_MC62"
TEMPLATE_PPTX = NB_DIR / "MC62_2Hz_staircase_presentation.pptx"
OUTPUT_PPTX = NB_DIR / "MC62_measurement_campaign.pptx"

# Executed notebooks -- in-place in the repository
EXEC_NBS = {
    "analysis_00": NB_DIR / "analysis" / "2026-02-11_00_test.ipynb",
    "analysis_01": NB_DIR / "analysis" / "2026-02-11_01_staircase_with_shims.ipynb",
    "analysis_02": NB_DIR / "analysis" / "2026-02-12_02_staircase_without_shims.ipynb",
    "analysis_03": NB_DIR / "analysis" / "2026-02-16_03_staircase_2Hz.ipynb",
    "analysis_04": NB_DIR / "analysis" / "2026-02-17_04_staircase_2Hz_morning.ipynb",
    "eddy_01": NB_DIR / "eddy_current" / "2026-02-11_01_staircase_with_shims.ipynb",
    "eddy_02": NB_DIR / "eddy_current" / "2026-02-12_02_staircase_without_shims.ipynb",
    "eddy_03": NB_DIR / "eddy_current" / "2026-02-16_03_staircase_2Hz.ipynb",
    "validation_02": NB_DIR / "validation" / "2026-02-12_02_vs_ffmm.ipynb",
    "compare_01v02": NB_DIR / "comparison" / "2026-02-12_01_vs_02_shims_effect.ipynb",
    "compare_03v04": NB_DIR / "comparison" / "2026-02-17_03_vs_04_reproducibility.ipynb",
}


# -- Image extraction ------------------------------------------------------

def extract_images(nb_path: Path) -> list[bytes]:
    """Extract all PNG images from an executed notebook, in cell order."""
    with open(nb_path, encoding="utf-8") as f:
        nb = json.load(f)

    images = []
    for cell in nb["cells"]:
        if cell["cell_type"] != "code":
            continue
        for output in cell.get("outputs", []):
            data = output.get("data", {})
            if "image/png" in data:
                b64 = data["image/png"]
                if isinstance(b64, list):
                    b64 = "".join(b64)
                images.append(base64.b64decode(b64))
    return images


def count_images_all():
    """Print image counts for all notebooks (use with --count)."""
    for key, path in EXEC_NBS.items():
        if path.exists():
            imgs = extract_images(path)
            print(f"  {key}: {len(imgs)} images")
        else:
            print(f"  {key}: NOT FOUND")


# -- PPTX helpers ----------------------------------------------------------

# Slide dimensions (CERN template standard: 13.33" x 7.5")
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Layout indices (from template inspection)
LAYOUT_TITLE = 1       # 'Title Slide'
LAYOUT_CONTENT = 3     # 'Content'
LAYOUT_CHAPTER = 17    # 'Chapter Header'
LAYOUT_TITLE_ONLY = 18 # 'Title Only'
LAYOUT_BLANK = 19      # 'Blank'
LAYOUT_LAST = 20       # 'Last slide'


def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    layout = prs.slide_layouts[LAYOUT_TITLE]
    slide = prs.slides.add_slide(layout)
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            shape.text = title
        elif shape.placeholder_format.idx == 1:
            shape.text = subtitle
    return slide


def add_chapter_slide(prs, chapter_num, title, subtitle=""):
    """Add a chapter header slide."""
    layout = prs.slide_layouts[LAYOUT_CHAPTER]
    slide = prs.slides.add_slide(layout)
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            shape.text = f"{chapter_num}. {title}"
        elif shape.placeholder_format.idx == 1:
            shape.text = subtitle
    return slide


def add_image_slide(prs, title, image_bytes, subtitle=""):
    """Add a slide with a title and a large image."""
    layout = prs.slide_layouts[LAYOUT_TITLE_ONLY]
    slide = prs.slides.add_slide(layout)

    # Set title
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            shape.text = title

    # Add image -- fill most of the slide below the title
    img_stream = io.BytesIO(image_bytes)
    max_w = Inches(12.7)
    max_h = Inches(5.8)
    top = Inches(1.4)

    # Get image aspect ratio
    from PIL import Image
    img = Image.open(io.BytesIO(image_bytes))
    img_w, img_h = img.size
    aspect = img_w / img_h

    # Fit within bounds
    if max_w / max_h > aspect:
        height = max_h
        width = int(height * aspect)
    else:
        width = max_w
        height = int(width / aspect)

    # Center horizontally
    left = (SLIDE_W - width) // 2

    slide.shapes.add_picture(img_stream, left, top, width, height)

    # Add subtitle if provided
    if subtitle:
        txBox = slide.shapes.add_textbox(
            Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
        )
        tf = txBox.text_frame
        tf.text = subtitle
        for p in tf.paragraphs:
            p.font.size = Pt(10)
            p.font.italic = True
            p.alignment = PP_ALIGN.LEFT

    return slide


def add_text_slide(prs, title, body_text):
    """Add a content slide with title and body text."""
    layout = prs.slide_layouts[LAYOUT_CONTENT]
    slide = prs.slides.add_slide(layout)
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            shape.text = title
        elif shape.placeholder_format.idx == 1:
            shape.text = body_text
    return slide


def add_last_slide(prs):
    """Add the CERN 'last slide'."""
    layout = prs.slide_layouts[LAYOUT_LAST]
    prs.slides.add_slide(layout)


def _safe(imgs, idx):
    """Return image at index if it exists, else None."""
    return imgs[idx] if len(imgs) > idx else None


# -- Main build ------------------------------------------------------------

def build():
    """Build the complete presentation."""
    # Load template (use existing PPTX for the master/layouts)
    prs = Presentation(str(TEMPLATE_PPTX))

    # Remove all existing slides
    from lxml import etree
    while len(prs.slides._sldIdLst) > 0:
        sldId = prs.slides._sldIdLst[0]
        rId = sldId.get(etree.QName(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships", "id"))
        if rId:
            prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(sldId)

    # Extract images from executed notebooks
    print("Extracting images from executed notebooks...")
    images = {}
    for key, path in EXEC_NBS.items():
        if path.exists():
            imgs = extract_images(path)
            images[key] = imgs
            print(f"  {key}: {len(imgs)} images")
        else:
            images[key] = []
            print(f"  {key}: NOT FOUND at {path}")

    # ==================================================================
    # Title Slide
    # ==================================================================
    add_title_slide(
        prs,
        "MC62 Rotating Coil Measurement Campaign",
        "LEAR C-shaped Dipole\n"
        "Staircase Tests -- Feb 11\u201317, 2026\n"
        "A. Bellelli -- CERN",
    )

    # ==================================================================
    # Setup & Pipeline Choices
    # ==================================================================
    add_text_slide(prs, "Measurement Setup & Analysis Pipeline",
        "5 staircase tests: 00 (check), 01 (shims, 1 Hz), "
        "02 (no shims, 1 Hz), 03 (2 Hz PM), 04 (2 Hz AM)\n\n"
        "Pipeline: dri + rot (cel/fed auto-disabled -- UNSAFE for dipole)\n"
        "R_ref = 33.0 mm, legacy drift mode, external Kn calibration\n"
        "Averaging: N_LAST = 170 (1 Hz) / 340 (2 Hz)\n\n"
        "PCBs: Integral (R45, 30 harmonics) + Central (DQ, 15 harmonics)\n"
        "Current cycle: 0\u2192+200\u21920\u2192\u2212200\u21920 A, "
        "20 A steps, 1 A/s ramp")

    # ==================================================================
    # Chapter 1: Test 00 -- System Check (Feb 11)
    # ==================================================================
    add_chapter_slide(prs, 1,
        "Test 00\nFeb 11, 2026 -- System Check",
        "Quick verification: 9 plateaus x 10 turns\n"
        "Short staircase: 0\u2192+200\u21920\u2192\u2212200\u21920 A\n"
        "With shims, 1 Hz rotation")

    # analysis_00 (7 images):
    # 0: current profile, 1: settling
    # 2: B1 hyst, 3: b2, 4: b3, 5: TF
    # 6: zoomed 2x2
    a00 = images.get("analysis_00", [])
    if _safe(a00, 6):
        add_image_slide(prs, "Test 00 -- Harmonic Hysteresis (I != 0)",
                        a00[6],
                        "System check: 10 turns/plateau, eddy currents not resolved")
    if _safe(a00, 2):
        add_image_slide(prs, "Test 00 -- B1 vs Current",
                        a00[2],
                        "b2 ~ \u2212152 units (C-shape asymmetry detected)")

    # ==================================================================
    # Chapter 2: Test 01 -- With Shims (Feb 11)
    # ==================================================================
    add_chapter_slide(prs, 2,
        "Test 01\nFeb 11, 2026 -- With Shims",
        "1 Hz staircase, 0\u2192+200\u21920\u2192\u2212200\u21920 A in 20 A steps\n"
        "41 plateaus, 350 turns each at \u221260 rpm\n"
        "Integral (R45) + Central (DQ) PCBs")

    # analysis_01 (run-based, 9 images):
    # 0: current profile, 1: settling visualization
    # 2: B1 hysteresis, 3: b2 hysteresis, 4: b3 hysteresis, 5: TF hysteresis
    # 6: full individual hyst, 7: zoomed 2x2 (I!=0), 8: multipole spectrum
    a01 = images.get("analysis_01", [])
    if _safe(a01, 7):
        add_image_slide(prs, "Test 01 -- Harmonic Hysteresis (I != 0)",
                        a01[7],
                        "B1, b2, b3, TF vs current -- b2 ~ \u2212151 units (with shims)")
    if _safe(a01, 2):
        add_image_slide(prs, "Test 01 -- B1 vs Current (Full Range)",
                        a01[2])
    if _safe(a01, 8):
        add_image_slide(prs, "Test 01 -- Multipole Spectrum at Peak Current",
                        a01[8])

    # Eddy current -- Test 01
    # eddy_01 (6 images): 0: settling, 1: fits, 2: tau vs I,
    # 3: bias (3-panel), 4: sensitivity, 5: double-exp
    e01 = images.get("eddy_01", [])
    if _safe(e01, 2):
        add_image_slide(prs, "Test 01 -- Eddy Current: \u03c4 vs Current",
                        e01[2],
                        "\u03c4 = 26.4 \u00b1 8.5 s (mean), "
                        "\u03bc_r dependence: 33 s (low I) \u2192 12 s (saturation)")
    if _safe(e01, 4):
        add_image_slide(prs, "Test 01 -- N_LAST_TURNS Sensitivity",
                        e01[4])

    # ==================================================================
    # Chapter 3: Test 02 -- Without Shims (Feb 12)
    # ==================================================================
    add_chapter_slide(prs, 3,
        "Test 02\nFeb 12, 2026 -- Without Shims",
        "1 Hz staircase, same cycle as Test 01\n"
        "Shims removed between Test 01 and Test 02")

    # analysis_02: same layout as analysis_01 (9 images)
    a02 = images.get("analysis_02", [])
    if _safe(a02, 7):
        add_image_slide(prs, "Test 02 -- Harmonic Hysteresis (I != 0)",
                        a02[7],
                        "B1, b2, b3, TF vs current -- b2 ~ \u221215 units (no shims!)")
    if _safe(a02, 2):
        add_image_slide(prs, "Test 02 -- B1 vs Current (Full Range)",
                        a02[2])
    if _safe(a02, 8):
        add_image_slide(prs, "Test 02 -- Multipole Spectrum at Peak Current",
                        a02[8])

    # Eddy current -- Test 02
    e02 = images.get("eddy_02", [])
    if _safe(e02, 2):
        add_image_slide(prs, "Test 02 -- Eddy Current: \u03c4 vs Current",
                        e02[2],
                        "\u03c4 = 26.0 \u00b1 8.1 s -- matches Test 01 "
                        "(shims do not affect eddy-current dynamics)")
    if _safe(e02, 4):
        add_image_slide(prs, "Test 02 -- N_LAST_TURNS Sensitivity",
                        e02[4])

    # ==================================================================
    # Chapter 4: Shims Effect -- 01 vs 02
    # ==================================================================
    add_chapter_slide(prs, 4,
        "Shims Effect\nTest 01 vs Test 02",
        "Quantify the effect of removing iron shims\n"
        "Blue = Test 01 (with shims), Red = Test 02 (without shims)\n"
        "Dark shade = ascending, Light shade = descending")

    # compare_01v02 (6 images):
    # 0: B1 overlay, 1: b2 overlay, 2: b3 overlay, 3: TF overlay
    # 4: difference bar charts (2x2), 5: multipole spectrum
    c12 = images.get("compare_01v02", [])
    if _safe(c12, 0):
        add_image_slide(prs, "Shims Effect -- B1 vs Current",
                        c12[0])
    if _safe(c12, 1):
        add_image_slide(prs, "Shims Effect -- b2 (Quadrupole) vs Current",
                        c12[1],
                        "b2 ~ \u2212151 with shims vs \u221215 without: "
                        "shims WORSENED the quadrupole")
    if _safe(c12, 2):
        add_image_slide(prs, "Shims Effect -- b3 (Sextupole) vs Current",
                        c12[2])
    if _safe(c12, 3):
        add_image_slide(prs, "Shims Effect -- Transfer Function vs Current",
                        c12[3])
    if _safe(c12, 4):
        add_image_slide(prs, "Shims Effect -- Per-Level Differences",
                        c12[4],
                        "\u0394 = Test 02 (no shims) \u2212 Test 01 (with shims)")
    if _safe(c12, 5):
        add_image_slide(prs, "Shims Effect -- Multipole Spectrum Comparison",
                        c12[5])

    add_text_slide(prs, "Shims Effect -- Key Findings",
        "Comparing Test 01 (with shims) vs Test 02 (without shims):\n\n"
        "\u2022 b2 shift: ~133 units -- shims INCREASED |b2| "
        "from 15 to 151 units\n"
        "\u2022 b2 correlation r = \u22120.515 -- field symmetry inverted\n"
        "\u2022 b3 shift: ~15 units (secondary, saturation redistribution)\n"
        "\u2022 B1 shift: ~5 mT (small -- shims affect homogeneity, "
        "not total flux)\n"
        "\u2022 TF shift: ~0.06 T/kA (~5%)\n\n"
        "\u2192 No-shims configuration has better field quality.\n"
        "   Shims need repositioning or removal.")

    # ==================================================================
    # Chapter 5: Validation -- Python vs FFMM C++ (Test 02)
    # ==================================================================
    add_chapter_slide(prs, 5,
        "Pipeline Validation\nPython vs FFMM C++ (Test 02)",
        "Run-based comparison: 41 plateaus x 350 turns\n"
        "FFMM options: dri rot nor cel fed dit\n"
        "Our options: dri rot cel fed (dit N/A on plateaus)")

    # validation_02 (5 images):
    # 0: N_LAST sweep, 1: parity scatter, 2: per-harmonic bars,
    # 3: per-run residuals, 4: Central PCB comparison
    v02 = images.get("validation_02", [])
    if _safe(v02, 0):
        add_image_slide(prs, "Validation -- N_LAST Averaging Window Sweep",
                        v02[0],
                        "FFMM uses all 350 turns. "
                        "Our 170-turn default differs by 72 \u00b5T (eddy settling)")
    if _safe(v02, 1):
        add_image_slide(prs, "Validation -- B_main Parity (N_LAST=350)",
                        v02[1],
                        "RMS = 0.6 \u00b5T -- sub-microtesla agreement")
    if _safe(v02, 2):
        add_image_slide(prs, "Validation -- Per-Harmonic Residuals",
                        v02[2],
                        "All harmonics < 0.003 units -- machine-precision parity")
    if _safe(v02, 3):
        add_image_slide(prs, "Validation -- Per-Run B_main Residuals",
                        v02[3])

    # ==================================================================
    # Chapter 6: Test 03 -- 2 Hz Staircase (Feb 16 afternoon)
    # ==================================================================
    add_chapter_slide(prs, 6,
        "Test 03\nFeb 16, 2026 -- Afternoon",
        "2 Hz staircase (120 rpm), 512 samples/turn\n"
        "~740 turns per plateau, streaming binary format\n"
        "No shims, same current cycle")

    # analysis_03 (streaming, 13 images):
    # 0: per-turn statistics, 1: annotated current profile
    # 2: turn classification map, 3: harmonics vs turn (with classification)
    # 4: all-turns field-vs-time (plot 1), 5: all-turns (plot 2)
    # 6: B1 hyst, 7: b2 hyst, 8: b3 hyst, 9: TF hyst
    # 10: zoomed 2x2 (I!=0), 11: multipole spectrum, 12: FFMM parity
    a03 = images.get("analysis_03", [])
    if _safe(a03, 1):
        add_image_slide(prs, "Test 03 -- Current Profile & Plateau Detection",
                        a03[1])
    if _safe(a03, 2):
        add_image_slide(prs, "Test 03 -- Turn Classification Map",
                        a03[2],
                        "Green = plateau, Orange = ramp, Grey = precycle")
    if _safe(a03, 10):
        add_image_slide(prs, "Test 03 -- Harmonic Hysteresis (I != 0)",
                        a03[10],
                        "B1, b2, b3, TF vs current -- b2 ~ \u221216, b3 ~ \u221212 units")
    if _safe(a03, 6):
        add_image_slide(prs, "Test 03 -- B1 vs Current (Full Range)",
                        a03[6])
    if _safe(a03, 11):
        add_image_slide(prs, "Test 03 -- Multipole Spectrum at Peak Current",
                        a03[11])
    if _safe(a03, 3):
        add_image_slide(prs, "Test 03 -- Harmonics vs Turn Number",
                        a03[3],
                        "B1, b2, b3 per turn with classification overlay")
    if _safe(a03, 4):
        add_image_slide(prs, "Test 03 -- All Turns: Field vs Time",
                        a03[4],
                        "Including ramp turns -- full current-harmonic correlation")
    if _safe(a03, 12):
        add_image_slide(prs, "Test 03 -- FFMM C++ Parity Check",
                        a03[12],
                        "B_main max |diff| = 1.81e-13 T (machine precision)")

    # Eddy current results for test 03
    e03 = images.get("eddy_03", [])
    # eddy_03 (7 images): 0: settling curves, 1: representative fits,
    # 2: tau vs current, 3: settling bias (3-panel),
    # 4: sensitivity study, 5: precycle, 6: double-exp comparison
    if _safe(e03, 0):
        add_image_slide(prs, "Test 03 -- Eddy Current: Settling Curves",
                        e03[0])
    if _safe(e03, 1):
        add_image_slide(prs, "Test 03 -- Eddy Current: Representative Fits",
                        e03[1])
    if _safe(e03, 2):
        add_image_slide(prs, "Test 03 -- Eddy Current: \u03c4 vs Current",
                        e03[2],
                        "\u03c4 = 13.2 \u00b1 8.7 s; "
                        "\u03bc_r dependence: 17 s (low I) \u2192 5 s (saturation)")
    if _safe(e03, 3):
        add_image_slide(prs, "Test 03 -- Settling Bias Analysis",
                        e03[3])
    if _safe(e03, 4):
        add_image_slide(prs, "Test 03 -- N_LAST_TURNS Sensitivity Study",
                        e03[4],
                        "b3 bias < 0.02 units for all N_LAST values")

    # ==================================================================
    # Chapter 7: Test 04 -- 2 Hz Staircase (Feb 17 morning)
    # ==================================================================
    add_chapter_slide(prs, 7,
        "Test 04\nFeb 17, 2026 -- Morning",
        "2 Hz staircase, repeat of Test 03\n"
        "Morning measurement for reproducibility check (~16 h apart)\n"
        "N_SKIP_END = 20 (trim last 20 turns)")

    # analysis_04 (streaming, 13 images): same layout as analysis_03
    a04 = images.get("analysis_04", [])
    if _safe(a04, 1):
        add_image_slide(prs, "Test 04 -- Current Profile & Plateau Detection",
                        a04[1])
    if _safe(a04, 2):
        add_image_slide(prs, "Test 04 -- Turn Classification Map",
                        a04[2],
                        "Green = plateau, Orange = ramp "
                        "(no precycle -- systematic stepping from start)")
    if _safe(a04, 10):
        add_image_slide(prs, "Test 04 -- Harmonic Hysteresis (I != 0)",
                        a04[10],
                        "B1, b2, b3, TF -- closely matches Test 03")
    if _safe(a04, 6):
        add_image_slide(prs, "Test 04 -- B1 vs Current (Full Range)",
                        a04[6])
    if _safe(a04, 11):
        add_image_slide(prs, "Test 04 -- Multipole Spectrum at Peak Current",
                        a04[11])
    if _safe(a04, 3):
        add_image_slide(prs, "Test 04 -- Harmonics vs Turn Number",
                        a04[3])
    if _safe(a04, 4):
        add_image_slide(prs, "Test 04 -- All Turns: Field vs Time",
                        a04[4])
    if _safe(a04, 12):
        add_image_slide(prs, "Test 04 -- FFMM C++ Parity Check",
                        a04[12],
                        "B_main max |diff| = 2.25e-12 T (machine precision)")

    # ==================================================================
    # Chapter 8: Reproducibility -- 03 vs 04
    # ==================================================================
    add_chapter_slide(prs, 8,
        "Reproducibility\nTest 03 vs Test 04",
        "Day-to-day: Feb 16 afternoon vs Feb 17 morning\n"
        "Identical hardware, same cycle, same settings\n"
        "Blue = Test 03, Red = Test 04")

    # compare_03v04 (8 images):
    # 0: B1 overlay, 1: b2 overlay, 2: b3 overlay, 3: TF overlay
    # 4: difference bars, 5: multipole spectrum, 6: hysteresis width
    # 7: B1 noise scatter
    c34 = images.get("compare_03v04", [])
    if _safe(c34, 0):
        add_image_slide(prs, "Reproducibility -- B1 vs Current",
                        c34[0])
    if _safe(c34, 1):
        add_image_slide(prs, "Reproducibility -- b2 vs Current",
                        c34[1])
    if _safe(c34, 2):
        add_image_slide(prs, "Reproducibility -- b3 vs Current",
                        c34[2])
    if _safe(c34, 3):
        add_image_slide(prs, "Reproducibility -- Transfer Function vs Current",
                        c34[3])
    if _safe(c34, 4):
        add_image_slide(prs, "Reproducibility -- Per-Level Differences",
                        c34[4],
                        "\u0394B1, \u0394b2, \u0394b3 (Test 04 \u2212 Test 03)")
    if _safe(c34, 5):
        add_image_slide(prs, "Reproducibility -- Multipole Spectrum Comparison",
                        c34[5])
    if _safe(c34, 6):
        add_image_slide(prs, "Reproducibility -- Hysteresis Width",
                        c34[6],
                        "Width reproducible to ~40 \u00b5T")
    if _safe(c34, 7):
        add_image_slide(prs, "Reproducibility -- Turn-to-Turn B1 Scatter",
                        c34[7],
                        "B1 std at each plateau -- measurement noise comparison")

    add_text_slide(prs, "Reproducibility -- Summary",
        "38 matched current levels at |I| > 0:\n\n"
        "\u2022 \u0394B1: max 62 \u00b5T (0.03% relative) -- "
        "ambient temperature drift\n"
        "\u2022 \u0394b2: max 0.27 units -- well within scatter\n"
        "\u2022 \u0394b3: max 0.05 units -- negligible\n"
        "\u2022 \u0394TF: max 0.002 T/kA\n\n"
        "Correlation: B1 r=1.0000, b2 r=0.9954, b3 r=0.9989\n\n"
        "\u2192 Excellent day-to-day reproducibility confirmed.\n"
        "   Measurement system is stable and reliable.")

    # ==================================================================
    # Chapter 9: Eddy Current Comparison -- All Tests
    # ==================================================================
    add_chapter_slide(prs, 9,
        "Eddy Current Analysis\nAll Tests Compared",
        "Single-exp: B1(t) = B_inf + A exp(\u2212t/\u03c4)\n"
        "Double-exp: B1(t) = B_inf + A\u2081 exp(\u2212t/\u03c4\u2081) "
        "+ A\u2082 exp(\u2212t/\u03c4\u2082)")

    if _safe(e01, 2):
        add_image_slide(prs, "Test 01 (With Shims) -- \u03c4 vs Current",
                        e01[2])
    if _safe(e02, 2):
        add_image_slide(prs, "Test 02 (No Shims) -- \u03c4 vs Current",
                        e02[2])

    # Double-exp comparison
    if _safe(e01, 5):
        add_image_slide(prs, "Test 01 -- Single vs Double Exp Fit",
                        e01[5])
    if _safe(e02, 5):
        add_image_slide(prs, "Test 02 -- Single vs Double Exp Fit",
                        e02[5])
    if _safe(e03, 6):
        add_image_slide(prs, "Test 03 (2 Hz) -- Single vs Double Exp Fit",
                        e03[6])

    add_text_slide(prs, "Eddy Current -- Key Findings",
        "1 Hz tests (01, 02): \u03c4 = 26 \u00b1 8 s\n"
        "2 Hz test (03): \u03c4 = 13 \u00b1 9 s\n\n"
        "\u2022 Clear \u03bc_r dependence: \u03c4 drops with |I| "
        "(higher permeability \u2192 lower time constant)\n"
        "\u2022 Shims do NOT affect eddy-current dynamics "
        "(01 vs 02: 26.4 vs 26.0 s)\n"
        "\u2022 b3 bias from eddy currents is negligible "
        "(< 0.03 units for all tests)\n"
        "\u2022 N_LAST_TURNS = 170 (1 Hz) / 340 (2 Hz) -- "
        "conservative, excludes 5\u00d7\u03c4_max\n"
        "\u2022 Double-exp: marginal R\u00b2 improvement -- "
        "single-exp adequate")

    # ==================================================================
    # Chapter 10: Summary & Conclusions
    # ==================================================================
    add_chapter_slide(prs, 10,
        "Summary\n& Conclusions", "")

    add_text_slide(prs, "MC62 Magnet Characterisation",
        "B1 at 200 A: 0.2185 T (integral)\n"
        "TF at 200 A: 1.093 T/kA\n"
        "Saturation onset: ~120 A\n\n"
        "b2 (no shims): \u221216 units (C-shape asymmetry)\n"
        "b2 (with shims): \u2212151 units (shims worsened it!)\n"
        "b3: \u221212 units (stable)\n\n"
        "Eddy current \u03c4: 2\u201340 s (current-dependent)\n"
        "Hysteresis width: 0.8\u20131.1 mT\n"
        "Reproducibility: \u0394B1 < 62 \u00b5T, \u0394b3 < 0.05 units")

    add_text_slide(prs, "Conclusions",
        "1. Shims effect: shims dramatically WORSENED b2 "
        "(151 vs 15 units). Need repositioning.\n\n"
        "2. Reproducibility: excellent (\u0394B1 \u2264 62 \u00b5T, "
        "\u0394b3 < 0.05 units) -- system is stable.\n\n"
        "3. Eddy currents: \u03c4 = 2\u201340 s with \u03bc_r dependence. "
        "N_LAST = 340 safely excludes transients.\n\n"
        "4. Pipeline validation: machine-precision parity with FFMM C++ "
        "(< 1 pT for B_main).\n\n"
        "5. cel/fed correctly auto-disabled: dipole high-order harmonics "
        "unreliable for centre-localisation.\n\n"
        "6. C-shape signature: systematic b2 ~ \u221216 units "
        "(inherent to open-gap geometry).")

    # -- Last slide --
    add_last_slide(prs)

    # Save
    prs.save(str(OUTPUT_PPTX))
    print(f"\nPresentation saved to: {OUTPUT_PPTX}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    if "--count" in sys.argv:
        print("Image counts per notebook:")
        count_images_all()
    else:
        build()
