"""Shared PPTX helper functions for CERN Corporate 16:9 template (2024).

All presentation scripts import these helpers instead of duplicating them.
Layout indices match the CERN Corporate template shipped in the repository.
"""

from __future__ import annotations

from pathlib import Path

import matplotlib.pyplot as plt
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


# ── CERN template layout indices ─────────────────────────────────────
LY_TITLE_LOGO  = 2   # "Title Slide with logo"
LY_CONTENT     = 3   # "Content"
LY_BULLET      = 4   # "Bulleted Content"
LY_BIG_PIC     = 5   # "Big Picture"
LY_TWO_CONTENT = 7   # "2 Contents"
LY_TXT_PIC     = 9   # "Text and Picture"
LY_CHAPTER     = 17  # "Chapter Header"
LY_TITLE_ONLY  = 18  # "Title Only"
LY_BLANK       = 19  # "1_Blank"
LY_LAST        = 20  # "Last slide"

# ── Table / text colours ─────────────────────────────────────────────
CERN_BLUE = RGBColor(0x00, 0x33, 0x99)
DARK_GREY = RGBColor(0x33, 0x33, 0x33)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xE8, 0xEB, 0xF0)

# Legacy aliases kept for scripts that import old names
HEADER_BG = CERN_BLUE
ROW_EVEN  = LIGHT_BG
ROW_ODD   = WHITE


# =====================================================================
#  Low-level helpers
# =====================================================================

def set_footer(slide, date=None, footer=None):
    """Set date (placeholder 10) and footer (placeholder 11) on *slide*.

    If *date* or *footer* is ``None`` the corresponding placeholder is
    left unchanged (i.e. whatever the template has).
    """
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 10 and date is not None:
            ph.text = date
        elif idx == 11 and footer is not None:
            ph.text = footer


def add_picture_to_slide(slide, img_path, max_w_in=11.8, max_h_in=5.3,
                         top_in=1.6, center=True):
    """Add a picture to *slide*, scaled to fit, centered horizontally."""
    from PIL import Image
    with Image.open(img_path) as im:
        iw, ih = im.size
    dpi = 180
    w_in = iw / dpi
    h_in = ih / dpi
    scale = min(max_w_in / w_in, max_h_in / h_in, 1.5)
    w, h = w_in * scale, h_in * scale
    left = (13.333 - w) / 2 if center else 0.5
    top = top_in + (max_h_in - h) / 2
    slide.shapes.add_picture(img_path, Inches(left), Inches(top),
                             Inches(w), Inches(h))


# =====================================================================
#  Slide builders
# =====================================================================

def slide_title_only(prs, title_text, img_path=None, notes="",
                     footer_date=None, footer_text=None):
    """Slide with just a title (layout 18) and optional image."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_ONLY])
    slide.placeholders[0].text = title_text
    set_footer(slide, date=footer_date, footer=footer_text)
    if img_path:
        add_picture_to_slide(slide, img_path)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def slide_chapter(prs, title_text, subtitle_text="",
                  footer_date=None, footer_text=None):
    """Section divider slide (layout 17)."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_CHAPTER])
    slide.placeholders[0].text = title_text
    if subtitle_text and 1 in slide.placeholders:
        slide.placeholders[1].text = subtitle_text
    set_footer(slide, date=footer_date, footer=footer_text)
    return slide


def slide_bullets(prs, title_text, bullets,
                  footer_date=None, footer_text=None):
    """Bulleted content slide (layout 4)."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_BULLET])
    slide.placeholders[0].text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if ": " in bullet:
            bold_part, rest = bullet.split(": ", 1)
            r1 = p.add_run()
            r1.text = bold_part + ": "
            r1.font.bold = True
            r1.font.size = Pt(16)
            r2 = p.add_run()
            r2.text = rest
            r2.font.size = Pt(16)
        else:
            r = p.add_run()
            r.text = bullet
            r.font.size = Pt(16)
        p.space_after = Pt(6)
    set_footer(slide, date=footer_date, footer=footer_text)
    return slide


def slide_text_and_pic(prs, title_text, bullets, img_path,
                       footer_date=None, footer_text=None):
    """Text on left, picture on right (layout 9)."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TXT_PIC])
    slide.placeholders[0].text = title_text
    # Left text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if ": " in bullet:
            bold_part, rest = bullet.split(": ", 1)
            r1 = p.add_run()
            r1.text = bold_part + ": "
            r1.font.bold = True
            r1.font.size = Pt(14)
            r2 = p.add_run()
            r2.text = rest
            r2.font.size = Pt(14)
        else:
            r = p.add_run()
            r.text = bullet
            r.font.size = Pt(14)
        p.space_after = Pt(4)
    # Right picture
    slide.placeholders[13].insert_picture(img_path)
    set_footer(slide, date=footer_date, footer=footer_text)
    return slide


def slide_table(prs, title_text, headers, rows, col_widths=None,
                footer_date=None, footer_text=None):
    """Content slide (layout 18 - Title Only) with a table below."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_ONLY])
    slide.placeholders[0].text = title_text
    set_footer(slide, date=footer_date, footer=footer_text)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    if col_widths is None:
        col_widths = [11.5 / n_cols] * n_cols
    table_w = sum(col_widths)
    left = (13.333 - table_w) / 2

    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(left), Inches(1.6),
        Inches(table_w), Inches(min(0.38 * n_rows, 5.5)))
    tbl = tbl_shape.table

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = CERN_BLUE

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = DARK_GREY
                p.alignment = PP_ALIGN.CENTER
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG
    return slide


def slide_two_images(prs, title_text, img_left, img_right,
                     label_left="", label_right="",
                     footer_date=None, footer_text=None):
    """Title-only slide with two images side by side."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_ONLY])
    slide.placeholders[0].text = title_text
    set_footer(slide, date=footer_date, footer=footer_text)
    from PIL import Image
    for img_path, x_off, label in [(img_left, 0.2, label_left),
                                   (img_right, 6.7, label_right)]:
        with Image.open(img_path) as im:
            iw, ih = im.size
        max_w, max_h = 6.2, 5.0
        scale = min(max_w / (iw / 180), max_h / (ih / 180))
        w, h = (iw / 180) * scale, (ih / 180) * scale
        top = 1.6 + (5.0 - h) / 2
        slide.shapes.add_picture(img_path, Inches(x_off), Inches(top),
                                 Inches(w), Inches(h))
        if label:
            txBox = slide.shapes.add_textbox(Inches(x_off), Inches(6.6),
                                             Inches(6.2), Inches(0.4))
            p = txBox.text_frame.paragraphs[0]
            p.text = label
            p.font.size = Pt(10)
            p.font.color.rgb = DARK_GREY
            p.alignment = PP_ALIGN.CENTER
    return slide


# =====================================================================
#  Figure save helper
# =====================================================================

def savefig(fig, output_dir, name):
    """Save *fig* as a PNG and close it.

    Parameters
    ----------
    fig : matplotlib Figure
    output_dir : str or Path
        Directory to write the PNG into.
    name : str
        Filename stem (without extension).

    Returns
    -------
    str
        Absolute path of the saved image.
    """
    path = Path(output_dir) / f"{name}.png"
    fig.savefig(str(path), bbox_inches="tight", dpi=180, facecolor="white")
    plt.close(fig)
    return str(path)
